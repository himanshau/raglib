"""Loader for Microsoft PowerPoint (.pptx) documents."""

from __future__ import annotations

import logging
import os
from typing import List

from raglib.schemas import Document

logger = logging.getLogger(__name__)


class PowerPointLoader:
    """Loads pptx files into Document objects."""

    def load(self, path: str) -> List[Document]:
        """Load a pptx path into a single Document."""

        try:
            from pptx import Presentation  # type: ignore[import-not-found]
        except ImportError as exc:
            raise ImportError(
                "PowerPoint loading requires python-pptx. Install with: pip install raglib-py[pptx]"
            ) from exc

        presentation = Presentation(path)
        lines: List[str] = []

        for slide_index, slide in enumerate(presentation.slides, start=1):
            lines.append(f"Slide {slide_index}")
            for shape in slide.shapes:
                text = getattr(shape, "text", "")
                if isinstance(text, str) and text.strip():
                    lines.append(text.strip())

        content = "\n".join(lines).strip()
        if not content:
            logger.warning("PowerPointLoader found no extractable text in file=%s", path)
            return []

        return [
            Document(
                id=f"pptx-{abs(hash(os.path.abspath(path)))}",
                content=content,
                metadata={"file": path, "loader": "powerpoint"},
                source="pptx",
            )
        ]
