"""PowerPoint loader for .pptx slide decks."""

from __future__ import annotations

import logging
import os
from typing import List

from raglib.schemas import Document

logger = logging.getLogger(__name__)


class PPTXLoader:
    """Load slide-level text from PowerPoint files into Document objects."""

    def load(self, path: str) -> List[Document]:
        """Load a .pptx file and return one Document per text-bearing slide."""

        try:
            from pptx import Presentation
        except ImportError as exc:
            raise ImportError(
                "PPTX loading requires python-pptx. Install with: pip install raglib-py[pptx]"
            ) from exc

        presentation = Presentation(path)
        base_hash = abs(hash(os.path.abspath(path)))
        documents: List[Document] = []

        for slide_index, slide in enumerate(presentation.slides, start=1):
            lines: List[str] = []
            for shape in slide.shapes:
                text = getattr(shape, "text", "")
                if text and text.strip():
                    lines.append(text.strip())

            content = "\n".join(lines).strip()
            if not content:
                continue

            documents.append(
                Document(
                    id=f"pptx-{base_hash}-slide-{slide_index}",
                    content=content,
                    metadata={"file": path, "slide": slide_index, "loader": "pptx"},
                    source="pptx",
                )
            )

        logger.info("PPTXLoader loaded %d documents from file=%s", len(documents), path)
        return documents
